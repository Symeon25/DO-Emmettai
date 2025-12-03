# file_loaders.py
"""
Shared loaders for all file types used in RAG.

- PDF / DOCX / TXT: LangChain's built-ins (local only, no HTTP)
- CSV: streaming CSVChunkedLoader
- Excel: streaming ExcelStreamingLoader
- PowerPoint: PptxRichLoader (per slide, with notes/tables)
"""

import os
from typing import List

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    Docx2txtLoader,
)

# ---------- Helpers for tabular & PPTX ----------

def _normalize_ws(s: str) -> str:
    """Normalize whitespace inside a string."""
    return " ".join(str(s).split()) if s is not None else ""


def _df_block_to_text(df) -> str:
    """Render a pandas DataFrame as pipe-separated text."""
    import pandas as pd

    if not isinstance(df, pd.DataFrame) or df.empty:
        return ""
    cols_line = " | ".join(map(str, df.columns.tolist()))
    lines = [cols_line]
    append = lines.append
    for _, row in df.iterrows():
        append(
            " | ".join(
                "" if pd.isna(v) else _normalize_ws(v)
                for v in row.tolist()
            )
        )
    return "\n".join(lines)


class ExcelStreamingLoader:
    """
    Streams ALL cells from ALL sheets using openpyxl read_only mode.

    - No header assumptions
    - All rows/cols included
    - Emits Documents per N rows to avoid huge single strings
    """

    def __init__(self, path: str, rows_per_chunk: int = 1000):
        self.path = path
        self.rows_per_chunk = rows_per_chunk

    @staticmethod
    def _safe_str(v):
        """Normalize any value to a clean string (handles None, dates, numbers, etc.)."""
        import datetime as _dt

        if v is None:
            return ""
        if isinstance(v, (_dt.datetime, _dt.date, _dt.time)):
            try:
                return v.isoformat()
            except Exception:
                return str(v)
        return str(v)

    def _rows_to_text(self, rows):
        """Serialize rows as pipe-separated lines (no header, no trimming)."""
        out = []
        append = out.append
        for r in rows:
            append(" | ".join(self._safe_str(c) for c in r))
        return "\n".join(out)

    def load(self) -> List[Document]:
        from openpyxl import load_workbook

        # read_only=True => streaming; data_only=True => values (not formulas)
        wb = load_workbook(self.path, read_only=True, data_only=True)
        docs: List[Document] = []

        for ws in wb.worksheets:
            rows_iter = ws.iter_rows(values_only=True)
            buffer, count, chunk_id, total_rows = [], 0, 0, 0

            for row in rows_iter:
                buffer.append(row)
                count += 1
                total_rows += 1

                if count >= self.rows_per_chunk:
                    text = (
                        f"[Excel sheet: {ws.title} | rows {total_rows - count + 1}-{total_rows}]\n"
                        + self._rows_to_text(buffer)
                    )
                    docs.append(
                        Document(
                            page_content=text,
                            metadata={
                                "source": self.path,
                                "sheet": ws.title,
                                "chunk_id": chunk_id,
                                "rows_in_chunk": count,
                                "rows_total_seen": total_rows,
                                "no_split": True,
                            },
                        )
                    )
                    buffer, count, chunk_id = [], 0, chunk_id + 1

            # Tail (remaining rows)
            if buffer:
                text = (
                    f"[Excel sheet: {ws.title} | rows {total_rows - count + 1}-{total_rows}]\n"
                    + self._rows_to_text(buffer)
                )
                docs.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": self.path,
                            "sheet": ws.title,
                            "chunk_id": chunk_id,
                            "rows_in_chunk": count,
                            "rows_total_seen": total_rows,
                            "no_split": True,
                        },
                    )
                )

            # Explicitly mark truly empty sheets
            if total_rows == 0:
                docs.append(
                    Document(
                        page_content=f"[Excel sheet: {ws.title}] (empty)",
                        metadata={
                            "source": self.path,
                            "sheet": ws.title,
                            "rows_total": 0,
                            "no_split": True,
                        },
                    )
                )

        return docs


class CSVChunkedLoader:
    """
    Reads ALL rows, ALL columns from CSV with streaming.
    Uses pandas read_csv(..., chunksize=...) to avoid memory blowups.
    Emits one Document per chunk—no trimming.
    """

    def __init__(self, path: str, chunksize: int = 20000, encoding: str | None = None):
        self.path = path
        self.chunksize = chunksize  # not a limit—just splits into many docs
        self.encoding = encoding

    def load(self) -> List[Document]:
        import pandas as pd

        docs: List[Document] = []
        chunk_iter = pd.read_csv(
            self.path,
            dtype=object,
            chunksize=self.chunksize,
            encoding=self.encoding,
            low_memory=False,
        )
        total_rows = 0
        chunk_id = 0
        for chunk in chunk_iter:
            rows = len(chunk)
            total_rows += rows
            text = (
                f"[CSV chunk {chunk_id} | rows {total_rows - rows}-{total_rows - 1}]\n"
                f"{_df_block_to_text(chunk)}"
            )
            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": self.path,
                        "chunk_id": chunk_id,
                        "rows_in_chunk": rows,
                        "rows_total_seen": total_rows,
                    },
                )
            )
            chunk_id += 1

        if not docs:  # empty file
            docs.append(
                Document(
                    page_content="[CSV] (empty file)",
                    metadata={"source": self.path, "rows_total_seen": 0},
                )
            )
        else:
            # annotate total rows on the last doc
            docs[-1].metadata["rows_total_final"] = total_rows

        return docs


def _table_to_text(table) -> str:
    """Extract text from python-pptx table."""
    rows = []
    for r in table.rows:
        cells = []
        for c in r.cells:
            cells.append(_normalize_ws(c.text))
        rows.append(" | ".join(cells))
    return "\n".join(rows)


class PptxRichLoader:
    """
    Extracts ALL textual content from slides (titles, shapes, tables, notes).
    Emits one Document per slide—no trimming.
    """

    def __init__(self, path: str):
        self.path = path

    def load(self) -> List[Document]:
        from pptx import Presentation

        prs = Presentation(self.path)
        docs: List[Document] = []
        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides, start=1):
            parts = []
            # Title placeholder
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    parts.append(f"[Title] {slide.shapes.title.text}")
            except Exception:
                pass

            # All shapes (text & tables)
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    txt = shape.text or ""
                    if txt.strip():
                        parts.append(txt)
                if hasattr(shape, "has_table") and shape.has_table:
                    parts.append("[Table]\n" + _table_to_text(shape.table))

            # Notes
            if (
                slide.has_notes_slide
                and slide.notes_slide
                and slide.notes_slide.notes_text_frame
            ):
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    parts.append("[Notes]\n" + notes)

            content = "\n".join(p for p in (s.strip() for s in parts) if p)
            docs.append(
                Document(
                    page_content=f"[Slide {idx}/{total_slides}]\n{content}",
                    metadata={
                        "source": self.path,
                        "slide_number": idx,
                        "slides_total": total_slides,
                    },
                )
            )

        if not docs:
            docs.append(
                Document(
                    page_content="[PowerPoint] (no readable text found)",
                    metadata={"source": self.path, "slides_total": 0},
                )
            )

        return docs


# ---------- Public entry point ----------

def pick_loader(path: str):
    """
    Decide which loader to use based on file extension.

    Policy:
    - Excel: all rows/cols (streamed)
    - CSV: all rows/cols (streamed)
    - PowerPoint: all text per slide
    - PDF/DOCX/TXT: LangChain loaders
    - Fallback: TextLoader
    """
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        return PyPDFLoader(path)

    if ext == ".docx":
        return Docx2txtLoader(path)

    if ext == ".txt":
        return TextLoader(path, autodetect_encoding=True)

    if ext in {".csv"}:
        # If you often see messy encodings, try encoding=None to let pandas detect.
        return CSVChunkedLoader(path, chunksize=20000, encoding=None)

    if ext in {".xlsx", ".xls"}:
        # Adjust rows_per_chunk for your environment; it does NOT drop data.
        return ExcelStreamingLoader(path, rows_per_chunk=800)

    if ext in {".pptx", ".ppt"}:
        return PptxRichLoader(path)

    # Fallback
    return TextLoader(path, autodetect_encoding=True)
